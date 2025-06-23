import tempfile
import base64
import os
import re
from pptx import Presentation
from unstructured.partition.pptx import partition_pptx
from unstructured.cleaners.core import clean, clean_extra_whitespace
from unstructured.cleaners.extract import  (
    extract_email_address,
    extract_ip_address,
    extract_us_phone_number,
    extract_datetimetz,
    extract_ordered_bullets
)

from rapidfuzz import fuzz



# # === CONFIG ===
# INPUT_PPTX = "example1.pptx"   # Replace with your test file
# OUTPUT_PPTX = "edited_output.pptx"
# COMMAND = "Change Q2 sales to 4500"

# # === STEP 1: Extract structure from the PPTX using Unstructured ===
# def extract_structure(pptx_path):
#     from collections import defaultdict
#     elements = partition_pptx(
#     filename=pptx_path,
#     include_metadata=True,
#     strategy="hi_res"
# )

#     slides = defaultdict(list)

#     for el in elements:
#         text = el.text.strip() if el.text else ""
#         if not text:
#             continue

#         meta = el.metadata.to_dict()
#         slide_num = int(meta.get("page_number", 0) or 0)

#         category = el.category or "Unknown"

#         slides[slide_num].append({
#             "category": category,
#             "text": text,
#             "metadata": {
#                 "slide_number": slide_num,
#                 "filename": meta.get("filename"),
#                 "coordinates": meta.get("coordinates"),  # Optional: for layout info
#                 "element_id": meta.get("element_id"),
#                 "emphasis": getattr(el, "emphasized_text_contents", None),
#             }
#         })

#     # Format as sorted list of slides
#     structured_output = []
#     for slide_num in sorted(slides):
#         structured_output.append({
#             "slide": slide_num,
#             "elements": slides[slide_num]
#         })

#     return structured_output


# # === STEP 2: Simulate parsing of the user's command (mock GPT) ===
# def parse_command(command: str):
#     match = re.search(r"change\s+(.*?)\s+to\s+([$\d,]+)", command, re.IGNORECASE)
#     if not match:
#         raise ValueError("Could not parse command.")
#     target = match.group(1).strip().lower()
#     value = match.group(2).replace(",", "").replace("$", "")
#     return target, value

# # === STEP 3: Find the matching table cell and update it ===
# def update_pptx(pptx_path, output_path, target_text, new_value):
#     prs = Presentation(pptx_path)
#     updated = False

#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if shape.has_table:
#                 table = shape.table
#                 for row in table.rows:
#                     for cell in row.cells:
#                         if fuzz.partial_ratio(cell.text.lower(), target_text) > 80:
#                             print(f"Updating '{cell.text}' to '{new_value}'")
#                             cell.text = new_value
#                             updated = True
#                             break
#                     if updated:
#                         break
#             if updated:
#                 break
#         if updated:
#             break

#     if updated:
#         prs.save(output_path)
#         print(f"✅ Saved updated presentation to {output_path}")
#     else:
#         print("⚠️ Could not find a matching table cell.")

# # === RUN IT ===
# if __name__ == "__main__":
#     print("🔍 Extracting content...")
#     structure = extract_structure(INPUT_PPTX)
#     import json

#     # Print all extracted elements clearly
#     print("\n📄 === Extracted Elements ===\n")
#     for slide in structure:
#         print(f"🔹 Slide {slide['slide']}")
#         for i, el in enumerate(slide["elements"], 1):
#             print(f"  {i}. [{el['category']}] {el['text']}")
#         print("-" * 50)
#     print("🧠 Slide content found:", len(structure), "elements")

#     print("📝 Parsing command:", COMMAND)
#     target, value = parse_command(COMMAND)

#     print(f"🎯 Looking for: '{target}' → '{value}'")
#     update_pptx(INPUT_PPTX, OUTPUT_PPTX, target, value)


'''
(!!!) THIS IS IT 
'''

from collections import defaultdict

# === CONFIG ===
INPUT_PPTX = "example1.pptx"  # Replace with your PPTX file path

# # === EXTRACT + CLEAN ===
# def extract_structure_cleaned(pptx_path):
#     elements = partition_pptx(
#         filename=pptx_path,
#         include_metadata=True,
#         strategy="hi_res",
#         infer_table_structure=True,
#     )

#     slides = defaultdict(list)

#     for el in elements:
#         raw_text = el.text or ""
#         cleaned_text = clean(
#             raw_text,
#             bullets=True,
#             extra_whitespace=True,
#             dashes=True,
#             trailing_punctuation=True,
#         )

#         if not cleaned_text:
#             continue

#         meta = el.metadata.to_dict()
#         slide_num = int(meta.get("page_number", 0) or 0)
#         category = el.category or getattr(el, "type", "Unknown")

#         slides[slide_num].append({
#             "category": category,
#             "cleaned_text": cleaned_text,
#             "original_text": raw_text,
#         })

#     return slides

def looks_like_table(text):
    lines = text.strip().splitlines()
    pipe_count = sum(1 for l in lines if '|' in l)
    tab_count = sum(1 for l in lines if '\t' in l)
    return len(lines) >= 2 and (pipe_count > 0 or tab_count > 0)

def label_element(text, category):
    if looks_like_table(text):
        return "LikelyTable"
    elif extract_email_address(text):
        return "Email"
    elif extract_ip_address(text):
        return "IPAddress"
    elif extract_us_phone_number(text):
        return "PhoneNumber"
    elif extract_datetimetz(text):
        return "DateTime"
    elif any(extract_ordered_bullets(text)):
        return "OrderedListItem"
    elif category == "UncategorizedText":
        return "Text"
    else:
        return category or "Unknown"

def extract_structure(pptx_path):
    elements = partition_pptx(filename=pptx_path, infer_table_structure=True)
    slides = defaultdict(list)

    for el in elements:
        raw_text = (el.text or "").strip()
        if not raw_text:
            continue

        meta = el.metadata.to_dict()
        slide_num = int(meta.get("page_number", 0) or 0)
        raw_category = el.category or "Unknown"

        labeled_type = label_element(raw_text, raw_category)

        slides[slide_num].append({
            "type": labeled_type,
            "text": raw_text,
            "metadata": {
                "slide_number": slide_num,
                "filename": meta.get("filename"),
                "coordinates": meta.get("coordinates"),
                "element_id": meta.get("element_id"),
                "raw_category": raw_category,
            }
        })

    structured_output = []
    for slide_num in sorted(slides):
        structured_output.append({
            "slide": slide_num,
            "elements": slides[slide_num]
        })

    return structured_output


# === DISPLAY RESULTS ===
def display_slides(slide_dict: list):
    print("\n📄 === Cleaned and Categorized Slide Content ===\n")
    for slide in sorted(slide_dict, key=lambda x: x["slide"]):
        print(f"🔹 Slide {slide['slide']}")
        for i, el in enumerate(slide["elements"], 1):
            print(f"  {i}. [{el['type']}] {el['text']}")
        print("-" * 50)


# === RUN ===
if __name__ == "__main__":
    print("🔍 Extracting and cleaning PPTX content...\n")
    structured_slides = extract_structure(INPUT_PPTX)
    display_slides(structured_slides)
